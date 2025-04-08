#!/usr/bin/python
# coding: utf-8


# <<<<<<<<<<<<<<<<<<<<<<<<<<      Imports        >>>>>>>>>>>>>>>>>>>>>>>>>>
import os
import re
import sys
import fitz  # pip install PyMuPDF
import json
import docx
import email
import time
import shutil
import hashlib
import exifread # pip install ExifRead
from email.policy import default
from striprtf.striprtf import rtf_to_text   # pip install striprtf

import argparse # for menu system
import msg_parser   # pip install msg_parser, extract-msg

from pptx import Presentation   # pip install python-pptx
# from pptx.oxml import nsmap # pip install nsmap

# import datetime
from datetime import datetime
from markdownify import markdownify as md   # pip install markdownify

# <<<<<<<<<<<<<<<<<<<<<<<<<<      Pre-Sets       >>>>>>>>>>>>>>>>>>>>>>>>>>

author = 'LincolnLandForensics'
description = "Convert the content of various file types to Markdown, for use in Obsidian"
version = '1.1.3'


global file_types
file_types = [
    '.bash', '.bat', '.bmp', '.c', '.cmd', '.cpp', '.cs', '.css', '.csv',
    '.docx', '.drawio', '.eml', '.flac', '.gif', '.go', '.heic', '.heif', '.htm', '.html',
    '.ini', '.java', '.jl', '.jpeg', '.jpg', '.js', '.json', '.kt', '.log',
    '.m', '.m4a', '.md', '.mermaid', '.mkv', '.mov', '.mp3', '.mp4', '.msg',
    '.ogg', '.ogv', '.pdf', '.php', '.png', '.ppt', '.pptx', '.ps1', '.py', '.r', '.rb',
    '.rs', '.rtf', '.sh', '.sql', '.svg', '.swift', '.tif', '.tiff', '.ts', '.tsv', '.txt',
    '.vbs', '.wav', '.webm', '.xlsx', '.xml', '.yaml', '.yml'
]

global files_docs
files_docs = [
    '.csv', '.docx', '.drawio', '.htm', '.html', '.json', '.md', '.mermaid', '.pdf', '.ppt', '.pptx',
    '.ps1', '.py', '.rtf', '.sh', '.tsv', '.txt', '.vbs', '.xlsx', '.xml',
    '.yaml', '.yml'
]

global plain_text
plain_text = [
    ".bat", ".cmd", ".csv", ".ini", ".json", ".log", ".md", 
    ".py", ".sh", ".ps1", ".txt", ".vbs", ".xml", ".yml", ".yaml"
]

global files_media
files_media = [
	'.bmp', '.flac', '.gif', '.heic', '.heif', '.jpeg', '.jpg', '.m4a', '.mkv', '.mov', '.mp3',
    '.mp4', '.ogg', '.ogv', '.png', '.svg', '.tiff', '.tif', '.wav', '.webm', '.webp'
]

global files_scripts
files_scripts = [
    '.bash', '.bat', '.c', '.cmd', '.go', '.java', '.js', '.m', '.php', '.ps1',
    '.py', '.r', '.rb', '.rs', '.sh', '.sql', '.swift', '.vbs', '.yaml', '.yml'
]


# Colorize section
global color_red
global color_yellow
global color_green
global color_blue
global color_purple
global color_reset
color_red = ''
color_yellow = ''
color_green = ''
color_blue = ''
color_purple = ''
color_reset = ''

if sys.version_info > (3, 7, 9) and os.name == "nt":
    version_info = os.sys.getwindowsversion()
    major_version = version_info.major
    build_version = version_info.build

    if major_version >= 10 and build_version >= 22000: # Windows 11 and above
        import colorama
        from colorama import Fore, Back, Style  
        print(f'{Back.BLACK}') # make sure background is black
        color_red = Fore.RED
        color_yellow = Fore.YELLOW
        color_green = Fore.GREEN
        color_blue = Fore.BLUE
        color_purple = Fore.MAGENTA
        color_reset = Style.RESET_ALL
        
# <<<<<<<<<<<<<<<<<<<<<<<<<<      Menu           >>>>>>>>>>>>>>>>>>>>>>>>>>

def main():
    """
    Main function to parse arguments and initiate file conversion.
    """
    parser = argparse.ArgumentParser(description=description)
    parser.add_argument('-I', '--input', help='Input folder path', required=False)
    parser.add_argument('-O', '--output', help='Output folder path', required=False)
    parser.add_argument('-T', '--Template', help='template path', required=False)

    parser.add_argument('-b', '--blank', help='create a blank obsidian folder', required=False, action='store_true')
    parser.add_argument('-c', '--convert', help='Convert files to markdown', required=False, action='store_true')
    parser.add_argument('-P', '--pdfconvert', help='markdown 2 pdf', required=False, action='store_true')
    parser.add_argument('-t', '--template', help='copy templated obsidian files/folders', required=False, action='store_true')

    args = parser.parse_args()

    global input_folder
    global output_folder
    global template1_input

    input_folder = os.path.abspath(args.input) if args.input else os.getcwd()
    output_folder = os.path.abspath(args.output) if args.output else r'ObsidianNotebook'
    template1_input = os.path.abspath(args.Template) if args.Template else 'Template_Cases'

    # Ensure the input folder exists
    if not os.path.exists(input_folder):

        msg_blurb = (f"Input folder {input_folder} doesn't exist.")
        msg_blurb_square(msg_blurb, color_red) 
        exit()  # test
        return 1
   
    if args.template:
        copy_template_folder(input_folder, output_folder, template1_input)
    if args.convert:
        msg_blurb = (f"reading input folder {input_folder}")
        msg_blurb_square(msg_blurb, color_green)         
        process_files(input_folder, output_folder)
    elif args.pdfconvert:
        read_markdown_files(input_folder,output_folder)
        msg_blurb = (f"PDF files saved to {output_folder}")
        msg_blurb_square(msg_blurb, color_green)        
    elif args.blank:
        copy_template_folder(input_folder, output_folder, template1_input)
        setup_obsidian(input_folder, output_folder)
    else:
        parser.print_help()
        Usage()

    return 0


# <<<<<<<<<<<<<<<<<<<<<<<<<<   Sub-Routines   >>>>>>>>>>>>>>>>>>>>>>>>>>

def cleanup_description(description):
    # Filter out keys with empty or None values, and exclude 'Icon' and 'Type'
    cleaned_description = {k: v for k, v in description.items() if v not in ('', None) and k not in ('Icon', 'Type', 'Latitude', 'Longitude', 'Time', 'FileTypeExtension', 'Name')}
    
    # Create a formatted string with key-value pairs
    formatted_description = "\n".join(f"{k}: {v}" for k, v in cleaned_description.items())
    
    return formatted_description


def copy_template_folder(input_folder, output_folder, template1_input):
    template1_input = os.path.join(input_folder, template1_input)
    # Check if the Templates folder exists in the input folder
    if os.path.exists(template1_input):
        # Copy contents from input Templates to output folder
        for item in os.listdir(template1_input):
            source_item = os.path.join(template1_input, item)
            destination_item = os.path.join(output_folder, item)
            
            if os.path.isdir(source_item):
                # If item is a directory and doesn't exist in destination, copy it recursively
                if not os.path.exists(destination_item):
                    shutil.copytree(source_item, destination_item, dirs_exist_ok=True)
                # else:
                    # print(f"Directory {destination_item} already exists.")
            else:
                # If item is a file and doesn't exist in destination, copy it
                if not os.path.exists(destination_item):
                    shutil.copy2(source_item, destination_item)
                else:
                    print(f"File {destination_item} already exists.")
        
        print(f"{template1_input} folder copied to {output_folder}")
    else:
        print("Templates folder does not exist in the input folder.")

def convert_to_decimal(coord, ref):
    """Converts GPS coordinates to decimal format."""
    if not coord:
        return None
    degrees = float(coord[0].num) / float(coord[0].den)
    minutes = float(coord[1].num) / float(coord[1].den) / 60.0
    seconds = float(coord[2].num) / float(coord[2].den) / 3600.0
    decimal = degrees + minutes + seconds
    if ref in ['S', 'W']:
        decimal = -decimal
    return decimal


def convert_to_markdown(file_path, output_folder):
    '''
    Function to convert the content of .txt, .pdf, .docx, and .eml to Markdown 
    ''' 
    file_name = os.path.basename(file_path)
    base_name, extension = os.path.splitext(file_name)
    target_file_path = os.path.join(output_folder, base_name + ".md")


    if file_path.lower().endswith(".pdf"):    
        content = extract_text_from_pdf(file_path) 
    elif file_path.lower().endswith(".pptx"):    
        content = extract_text_from_pptx(file_path) 
    elif extension.lower() in (files_media):
        content = md_for_media(file_path)

    elif extension.lower() in (".docx"):
        content = extract_text_from_docx_with_formatting(file_path)
    elif file_path.lower().endswith(".eml"):
        content = extract_text_from_eml(file_path)
    elif file_path.lower().endswith(".msg"):
        content = extract_text_from_msg(file_path)
    elif file_path.lower().endswith(".rtf"):
        content = convert_rtf(file_path)
    elif file_path.lower().endswith(('.html', '.htm')):
        content = extract_html(file_path)
    elif extension.lower() in (files_scripts):
        content = md_for_scripts(file_path)

    elif extension.lower() in plain_text:
        content = parse_text_file(file_path)

    else:
        print(f"{color_red}Unsupported file type: {color_reset}{file_path}")
        return
    
    # Write the content to a Markdown file
    with open(target_file_path, 'w', encoding='utf-8') as md_file:
        md_file.write(content)
    # try:
        # with open(target_file_path, 'w', encoding='utf-8') as md_file:
            # md_file.write(content)
    # except Exception as e:
        # return f"Error writing {file_path}: {e}"
    
    print(f"{color_green}Converted {color_reset}'{file_name}' to '{base_name}.md'.")


def convert_rtf(file_path):
    """
    Converts an RTF file to Markdown format by extracting the plain text.
    Formatting will be basic (like bold, italics, etc.), but more advanced formatting would need
    manual conversion.
    """
    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = "" 


    # Read the RTF file
    try:
        with open(file_path, 'r', encoding='utf-8') as rtf_file:
            rtf_content = rtf_file.read()
    except Exception as e:
        return f"Error reading RTF file: {e}"

    # Convert the RTF content to plain text
    plain_text = rtf_to_text(rtf_content)

    # Convert plain text into Markdown:
    text = convert_rtf_to_md2(plain_text)
    (creation_time, access_time, modified_time) = get_file_timestamps(file_path)
    file_name = os.path.basename(file_path)
    
    new_file_path = os.path.join(docs_folder, file_name)
    # Copy the file to the new path
    try:
        shutil.copy(file_path, new_file_path)
    except Exception as e:
        print(f"Error copying file '{file_path}': {e}")    
        
    new_file_path_link = os.path.join('4.Documents/', file_name) 
    new_file_path_link = new_file_path_link.replace("\\", "/")    
    
    
    # text = (f"\n## File: {file_path}\n## Creation: {creation_time}\n## Modified: {modified_time}\n\n{text}")                    
    text = (f"\n## [File]({new_file_path_link}): {file_path}\n"
                f"## Creation: {creation_time}\n"
                f"## Modified: {modified_time}\n"
                f"## MD5: {hash_md5}\n\n"                
                f"{text}\n\n"
                )        
    return text


def convert_rtf_to_md2(plain_text):
    """
    Converts the plain text extracted from the RTF to basic Markdown.
    This function is designed for simple formatting.
    """
    # Example: Convert any instances of **bold** text to Markdown
    markdown_content = plain_text
    markdown_content = re.sub(r'\\b (.*?) \\b0', r'**\1**', markdown_content)  # Bold text
    markdown_content = re.sub(r'\\i (.*?) \\i0', r'*\1*', markdown_content)  # Italics text
    markdown_content = re.sub(r'\\ul (.*?) \\ulnone', r'_\1_', markdown_content)  # Underline text
    markdown_content = re.sub(r'•\s*(.*?)\n', r'- \1\n', markdown_content)  # bullet lists
    markdown_content = re.sub(r'\\pard (.*?) \\pard0', r'\1', markdown_content)  # Paragraphs
    markdown_content = re.sub(r'\\pard\s*(.*?)\\pard0', r'\1\n\n', markdown_content)    # Convert paragraphs
    markdown_content = markdown_content.strip() # Clean up extra spaces and formatting errors

    try:
        markdown_content = re.sub(r'**(.*?)**\n', r'# \1\n', markdown_content)  # simple bold headings
    except:
        pass

    return markdown_content

def extract_text_from_docx_with_formatting(file_path, docs_folder='4.Documents/'):
    """
    Function to extract formatted text from a DOCX file, return it in Markdown format, 
    and include metadata at the end of the extracted text.
    """
    # Get file timestamps (creation, access, modified)
    creation_time, access_time, modified_time = get_file_timestamps(file_path)

    # Generate MD5 hash for the file
    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = ""  

    # Initialize the text variable
    text = ""

    try:
        # Copy the DOCX file to the assets folder
        file_name = os.path.basename(file_path)
        new_file_path_link = os.path.join(docs_folder, file_name).replace("\\", "/")
        new_file_path = os.path.join(docs_folder, file_name)

        # Copy the file to the new path
        try:
            shutil.copy(file_path, new_file_path)
        except Exception as e:
            print(f"Error copying file '{file_path}': {e}")
        # Initialize the text output with metadata and the file link
        text = (f"\n## [File]({new_file_path_link}): {file_path}\n"
                f"## Creation: {creation_time}\n"
                f"## Modified: {modified_time}\n"
                f"## MD5: {hash_md5}\n\n")

    except Exception as e:
        print(f"Error copying document '{file_path}': {e}")

    try:
        # Load DOCX file
        doc = docx.Document(file_path)

        # Extract formatted text from paragraphs
        for para in doc.paragraphs:
            para_text = handle_paragraph_formatting(para)
            text += para_text + "\n"

        # Extract metadata
        core_properties = doc.core_properties
        metadata = {
            "title": core_properties.title,
            "subject": core_properties.subject,
            "author": core_properties.author,
            "keywords": core_properties.keywords,
            "comments": core_properties.comments,
            "created": core_properties.created,
            "last_modified_by": core_properties.last_modified_by,
            "modified": core_properties.modified,
            "revision": core_properties.revision,
        }

        # Append metadata to the text
        text += "\n\n## Metadata\n"
        for key, value in metadata.items():
            if value:
                text += f"{key}: {value}\n"
            # else:
                # text += f"{key}: Not available\n"

    except Exception as e:
        print(f"Error reading DOCX file '{file_path}': {e}")

    # Return the extracted text with metadata and Markdown formatting
    return text
    


def extract_text_from_eml(file_path):
    """
    Function to extract text from EML file, including metadata and email content.
    """
    # Get file timestamps
    creation_time, access_time, modified_time = get_file_timestamps(file_path)
    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = ""  

    # Read the EML file
    with open(file_path, 'r', encoding='utf-8') as f:
        msg = email.message_from_file(f, policy=default)

    # Extract metadata
    text = (
        f"\n## File: {file_path}\n"
        f"## Creation: {creation_time}\n"
        f"## Modified: {modified_time}\n"
        f"## MD5: {hash_md5}\n\n"
        f"Subject: {msg['subject']}\n"
        f"From: {msg['from']}\n"
        f"To: {msg['to']}\n\n"
    )

    # Extract email body
    if msg.is_multipart():
        for part in msg.iter_parts():
            content_type = part.get_content_type()
            if content_type == "text/plain":
                text += part.get_content()
            elif content_type == "text/html":
                text += f"\n[HTML Content Skipped for Readability]\n"
    else:
        text += msg.get_content()

    return text


def extract_html(file_path):

    (creation_time, access_time, modified_time) = get_file_timestamps(file_path)    # test
    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = "" 
        
    try:
        # Copy the doc to the assets folder

        file_name = os.path.basename(file_path)

        new_file_path_link = os.path.join('9.Assets', file_name) 
        new_file_path_link = new_file_path_link.replace("\\", "/")
        # new_file_path = os.path.join(assets_folder, file_name)
        new_file_path = os.path.join(docs_folder, file_name)

        # Copy the file to the new path
        try:
            shutil.copy(file_path, new_file_path)
        except Exception as e:
            print(f"Error copying file '{file_path}': {e}")
        # Initialize the text output with metadata and the file link
        text = (f"\n## [File]({new_file_path_link}): {file_path}\n"
                f"## Creation: {creation_time}\n"
                f"## Modified: {modified_time}\n"
                f"## MD5: {hash_md5}\n\n"
                )
    except Exception as e:
        print(f"{color_red}Error copying document {color_reset}'{file_path}': {e}")


    # Read the HTML content from the file
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()

        # Convert the HTML content to Markdown using markdownify
        text = (f'({text}\n{md(html_content)}')
 
    except UnicodeDecodeError:
        print(f"{color_red}Cannot decode {color_reset} {file_path}.") 

    return text
  
  
def extract_text_from_msg(file_path):
    """
    Function to extract text from MSG files.
    """
    from extract_msg import Message

    # Get file timestamps
    creation_time, access_time, modified_time = get_file_timestamps(file_path)

    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = "" 
        
    # Parse the .msg file
    msg = Message(file_path)

    # Extract metadata and message content
    text = (
        f"\n## File: {file_path}\n"
        f"## Creation: {creation_time}\n"
        f"## Modified: {modified_time}\n"
        f"## MD5: {hash_md5}\n\n"        
        f"Subject: {msg.subject}\n"
        f"From: {msg.sender}\n"
        f"To: {msg.to}\n\n"
        f"Body:\n{msg.body}\n"
    )

    # Replace 'From: None' with 'From: '
    text = text.replace('From: None', 'From: ')
    
    return text


def extract_text_from_pdf(file_path):
    """
    Function to extract text from a PDF file, copy it to a new location in the assets folder,
    and generate a link to the original file. Appends metadata at the end under a "Metadata" heading.
    """

    if 1==1:
    # try:
        # Open the PDF file
        try:
            doc = fitz.open(file_path)
        except Exception as e:
            print(f"Error reading file '{file_path}': {e}")
            return ""
        
        # Get file timestamps
        creation_time, access_time, modified_time = get_file_timestamps(file_path)

        # If timestamps are not found, return empty text
        if creation_time is None or access_time is None or modified_time is None:
            return ""

        # Copy the PDF to the assets folder
        if not os.path.exists(assets_folder):
            os.makedirs(assets_folder)  # Create the assets folder if it doesn't exist

        file_name = os.path.basename(file_path)

        new_file_path_link = os.path.join('9.Assets/documents', file_name) 
        new_file_path_link = new_file_path_link.replace("\\", "/")
        new_file_path = os.path.join(docs_folder, file_name)

        # Copy the file to the new path
        try:
            shutil.copy(file_path, new_file_path)
        except Exception as e:
            print(f"Error copying file '{file_path}': {e}")
        # Initialize the text output with metadata and the file link
        text = (f"\n## [File]({new_file_path_link}): {file_path}\n"
                f"## Creation: {creation_time}\n"
                f"## Modified: {modified_time}\n\n")

        # Extract text from each page
        for page_num in range(doc.page_count):
            try:
                page = doc.load_page(page_num)
                page_text = page.get_text("text")  # Extract plain text
                text += page_text
            except Exception as e:
                print(f"An error occurred while processing the .pdf file: {e}")
            
        # Extract metadata
        metadata = doc.metadata

        if metadata:
            text += "\n\n## Metadata\n"
            for key, value in metadata.items():
                text += f"{key}: {value}\n"
        # else:
            # text += "\n\n## Metadata\nNo metadata found.\n"

        return text

    # except fitz.EmptyFileError:
        # print(f"Cannot open empty or corrupted PDF file: {file_path}")
        # return ""
    # except Exception as e:
        # print(f"Error reading PDF file '{file_path}': {e}")
        # return ""

def extract_text_from_pptx(file_path):
    """
    Extract text from a PowerPoint (.pptx) file, including metadata.
    Text is followed by a "Metadata" section with the file's properties.
    """
    text = ''
    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = ""

    try:
        # Get file timestamps
        creation_time, access_time, modified_time = get_file_timestamps(file_path)


        
        # Load the PowerPoint presentation
        presentation = Presentation(file_path)

        # Extract text from slides
        text = ""
        for slide_num, slide in enumerate(presentation.slides, start=1):
            text += f"\n## Slide {slide_num}:\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += paragraph.text + "\n"

        # Extract metadata using core_properties
        core_properties = presentation.core_properties
        metadata = {
            "title": core_properties.title,
            "subject": core_properties.subject,
            "author": core_properties.author,
            "keywords": core_properties.keywords,
            "comments": core_properties.comments,
            "created": core_properties.created,
            "last_modified_by": core_properties.last_modified_by,
            "modified": core_properties.modified,
            "revision": core_properties.revision,
        }

        # Append metadata to the text
        text += "\n\n## Metadata\n"
        for key, value in metadata.items():
            if value:
                text += f"{key}: {value}\n"

        # return text

    except Exception as e:
        print(f"An error occurred while processing the PPTX file: {e}")
        return ""

    file_name = os.path.basename(file_path)
    
    new_file_path = os.path.join(docs_folder, file_name)
    # Copy the file to the new path
    try:
        shutil.copy(file_path, new_file_path)
    except Exception as e:
        print(f"Error copying file '{file_path}': {e}")    
    new_file_path_link = os.path.join('4.Documents/', file_name) 
    new_file_path_link = new_file_path_link.replace("\\", "/")    
    
    
    # text = (f"\n## File: {file_path}\n## Creation: {creation_time}\n## Modified: {modified_time}\n\n{text}")                    
    text = (f"\n## [File]({new_file_path_link}): {file_path}\n"
                f"## Creation: {creation_time}\n"
                f"## Modified: {modified_time}\n"
                f"## MD5: {hash_md5}\n\n"                
                f"{text}\n\n"
                )    


    return text


def get_file_timestamps(file_path):
    """
    Returns the creation, last accessed, and last modified times of a file in YYYY-MM-DD hh:mm:ss format.

    :param file_path: The file path for which timestamps are needed.
    :return: A tuple containing (creation_time, access_time, modified_time)
    """
    if not os.path.isfile(file_path):
        raise ValueError(f"The path {file_path} is not a valid file.")

    # Get the file stats
    file_stats = os.stat(file_path)

    # Get the creation time (Windows), access time, and modification time
    creation_time = file_stats.st_ctime  # Windows: creation time, Linux: change time
    access_time = file_stats.st_atime
    modified_time = file_stats.st_mtime

    # Convert from timestamps to human-readable YYYY-MM-DD hh:mm:ss format
    creation_time = datetime.fromtimestamp(creation_time).strftime('%Y-%m-%d %H:%M:%S')
    access_time = datetime.fromtimestamp(access_time).strftime('%Y-%m-%d %H:%M:%S')
    modified_time = datetime.fromtimestamp(modified_time).strftime('%Y-%m-%d %H:%M:%S')

    return creation_time, access_time, modified_time


def handle_paragraph_formatting(para):
    '''
    Convert paragraph content to Markdown format considering styles.
    '''
    para_text = ""
    
    # Check if the paragraph is a heading (Heading 1, Heading 2, etc.)
    if para.style.name.startswith('Heading'):
        level = int(para.style.name.split()[-1])  # Extract level of the heading
        para_text += "#" * level + " " + para.text
    else:
        # For normal paragraphs, handle runs (bold, italic, etc.)
        for run in para.runs:
            text = run.text
            
            # Bold text
            if run.bold:
                text = f"**{text}**"
                
            # Italic text
            if run.italic:
                text = f"*{text}*"
            
            para_text += text

    return para_text    


def hashfile(file_path):
    """
    Computes and returns the MD5 hash of a file.
    """
    # Create an MD5 hash object
    hash_md5 = hashlib.md5()

    # Open the file in binary read mode
    with open(file_path, "rb") as f:
        # Read the file in chunks to avoid memory overload with large files
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    
    # Return the hexadecimal representation of the hash
    return hash_md5.hexdigest()
    
  
def parse_text_file(file_path):
    """
    Parses plain text files (.txt, .md, .cmd, .py), handles encoding errors,
    and returns the content as a string formatted for Markdown.
    """

    file_name = os.path.basename(file_path)

    creation_time, access_time, modified_time = get_file_timestamps(file_path)
    new_file_path_link = os.path.join('4.Documents/', file_name) 
    new_file_path_link = new_file_path_link.replace("\\", "/")

    new_file_path = os.path.join(docs_folder, file_name)
    
    # Copy the file to the new path
    try:
        shutil.copy(file_path, new_file_path)
    except Exception as e:
        print(f"Error copying file '{file_path}': {e}")
    header = f"\n## [File]({new_file_path_link}): {file_path}\n## Creation: {creation_time}\n## Modified: {modified_time}\n\n"

    header = f"\n## File: {file_path}\n## Creation: {creation_time}\n## Modified: {modified_time}\n\n"
    base_name, extension = os.path.splitext(file_path)
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='ISO-8859-1') as f:
                content = f.read()
        except UnicodeDecodeError:
            print(f"{color_red}Cannot decode {color_reset} {file_path} with UTF-8 or ISO-8859-1.") 
            # return None

    if extension.lower() in (".py"):
        return header + "\n\n\n" + "```python\n" +content + "```" + "\n\n\n"
        # Check if the .obsidian folder exists, if not, create it
        if not os.path.exists('scripts'):
            os.makedirs('scripts')
            print(f"scripts folder created at scripts")
  
    elif extension.lower() in (".ps1"):
        return header + "\n\n\n" + "```powershell\n" +content + "\n```" + "\n\n\n"
    elif extension.lower() in (".sh"):
        return header + "\n\n\n" + "```bash\n" +content + "\n```" + "\n\n\n"
    elif extension.lower() in (".cmd", ".bat"):
        return header + "\n\n\n" + "```cmd\n" +content + "\n```" + "\n\n\n"

    elif extension.lower() in (".vbs"):
        return header + "\n\n\n" + "```vbs\n" +content + "\n```" + "\n\n\n"
    else:
        return header + content + "\n"

    
def md_for_media(file_path):
    """
    copy media files to a new location in the media folder,
    and generate a link to the original file.
    """
    Description = ''
    try:
        # Get file timestamps
        creation_time, access_time, modified_time = get_file_timestamps(file_path)

        # If timestamps are not found, return empty text
        if creation_time is None or access_time is None or modified_time is None:
            print('error converting time stamp')
            # return ""

    except Exception as e:
        print(f"Error making file '{file_path}': {e}")
        
    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = ""

    file_name = os.path.basename(file_path)

    new_file_path_link = os.path.join('9.Assets/media/', file_name) 
    new_file_path_link = new_file_path_link.replace("\\", "/")

        
    new_file_path = os.path.join(media_folder, file_name)

    # Copy the file to the new path
    try:
        shutil.copy(file_path, new_file_path)
    except Exception as e:
        print(f"Error copying file '{file_path}': {e}")
    try:
        if file_path.lower().endswith(('.heic', '.heif', '.jpg', '.jpeg', '.png', '.tiff', '.tif', '.webp')):
            exif_data, Description = read_exif_data(file_path)

    except Exception as e:
        print(f"{color_red}Error getting exif data '{file_path}'{color_reset}: {e}")
        

    if Description != '':
        Description = (f'## Exif Data\n{Description}')

    # Initialize the text output with metadata and the file link
    text = (f"\n### [File]({new_file_path_link}): {file_path}\n"
        f"### Creation: {creation_time}\n"
        f"### Modified: {modified_time}\n"
        f"### MD5: {hash_md5}\n\n"                
        f"### ![[{new_file_path_link}|400]]\n\n"
        f"{Description}\n"
        )

    return text
 

def md_for_scripts(file_path):
    """
    copy scripts files to a new location in the scripts folder,
    and generate a link to the original file.
    """
    try:
        hash_md5 = hashfile(file_path)
    except:
        hash_md5 = "" 

    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='ISO-8859-1') as f:
                content = f.read()
        except UnicodeDecodeError:
            print(f"{color_red}Cannot decode {color_reset} {file_path} with UTF-8 or ISO-8859-1.") 
            # return None




    try:
        # Get file timestamps
        creation_time, access_time, modified_time = get_file_timestamps(file_path)

        # If timestamps are not found, return empty text
        if creation_time is None or access_time is None or modified_time is None:
            return ""


        base_name, extension = os.path.splitext(file_path)


        file_name = os.path.basename(file_path)

        new_file_path_link = os.path.join('9.Assets/scripts/', file_name) 
        # new_file_path_link = os.path.join(scripts_folder, file_name) 
        new_file_path_link = new_file_path_link.replace("\\", "/")
        
        new_file_path = os.path.join(scripts_folder, file_name)

        # Copy the file to the new path
        try:
            shutil.copy(file_path, new_file_path)
        except Exception as e:
            print(f"Error copying file '{file_path}': {e}")
        
        # don't put a link in, it will run the script
        # header = f"\n## [File]({new_file_path_link}): {file_path}\n## Creation: {creation_time}\n## Modified: {modified_time}\n\n"

        header = f"\n## File: {file_path}\n## Creation: {creation_time}\n## Modified: {modified_time}\n\n"

    
        # new 
        if extension.lower() in (".py"):
            return header + "\n\n\n" + "```python\n" +content + "```" + "\n\n\n"

        elif extension.lower() in (".ps1"):
            return header + "\n\n\n" + "```powershell\n" +content + "\n```" + "\n\n\n"
        elif extension.lower() in (".sh"):
            return header + "\n\n\n" + "```bash\n" +content + "\n```" + "\n\n\n"
        elif extension.lower() in (".cmd", ".bat"):
            return header + "\n\n\n" + "```cmd\n" +content + "\n```" + "\n\n\n"

        elif extension.lower() in (".vbs"):
            return header + "\n\n\n" + "```vbs\n" +content + "\n```" + "\n\n\n"
        else:
            return header + content + "\n"


        # Initialize the text output with metadata and the file link
        text = (f"\n## [File]({new_file_path_link}): {file_path}\n"
                f"## Creation: {creation_time}\n"
                f"## Modified: {modified_time}\n"
                f"## MD5: {hash_md5}\n\n"
                f"## header test: \n{header}\n\n"
                
                f"## ![[{new_file_path_link}]]\n\n")

        return text
    except Exception as e:
        print(f"Error making file '{file_path}': {e}")
        return ""    

def md2pdf(filename, input_folder, output_folder):
    today_str = datetime.today().strftime('%Y-%m-%d')
    file_path = os.path.join(input_folder, filename)

    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            lines = file.readlines()
    except Exception as e:
        print(f"Failed to read {file_path}: {e}")
        return

    # Filter lines
    filtered_lines = [
        line.strip() for line in lines
        if line.strip() and not any(
            keyword in line.lower()
            for keyword in ['kanban', '```', '%%', '---', '**complete**']
        )
    ]

    # Add header
    header = [f"Filename: {filename}", f"Date: {today_str}", ""]
    all_lines = header + filtered_lines

    # Save to .txt if more than 50 lines
    if len(all_lines) > 50:
        txt_filename = os.path.join(output_folder, f"{filename.replace('.md', '')}_{today_str}.txt")
        try:
            with open(txt_filename, 'w', encoding='utf-8') as txt_file:
                txt_file.write('\n'.join(all_lines))
            print(f"\t{txt_filename}")
        except Exception as e:
            print(f"Failed to write TXT {txt_filename}: {e}")

    # Split into 50-line chunks
    chunk_size = 50
    chunks = [all_lines[i:i + chunk_size] for i in range(0, len(all_lines), chunk_size)]

    # Create PDF
    pdf_document = fitz.open()

    for index, chunk in enumerate(chunks):
        text = '\n'.join(chunk).strip()
        if not text:
            continue  # Skip empty pages

        page = pdf_document.new_page()
        text_rect = fitz.Rect(50, 50, 550, 800)
        inserted_chars = page.insert_textbox(text_rect, text, fontsize=12, fontname="helv")

        if inserted_chars == 0:
            print(f"Warning: No text inserted on page {index + 1} for {filename}")

    pdf_filename = os.path.join(output_folder, f"{filename.replace('.md', '')}_{today_str}.pdf")

    try:
        pdf_document.save(pdf_filename)
        print(f"\t{pdf_filename}")
    except Exception as e:
        print(f"Failed to save PDF {pdf_filename}: {e}")
    finally:
        pdf_document.close()


def msg_blurb_square(msg_blurb, color):
    horizontal_line = f"+{'-' * (len(msg_blurb) + 2)}+"
    empty_line = f"| {' ' * (len(msg_blurb))} |"

    print(color + horizontal_line)
    print(empty_line)
    print(f"| {msg_blurb} |")
    print(empty_line)
    print(horizontal_line)
    print(f'{color_reset}')
    
    
def process_files(root_folder, output_folder):
    '''
    Function to walk through all files and convert
    ''' 
    
    setup_obsidian(input_folder, output_folder)   # create a default obsidian setup, if it doesn't exist

    # Walk through all files and subdirectories in root_folder
    for subdir, _, files in os.walk(root_folder):
        for file in files:
            file_path = os.path.join(subdir, file)
            base_name, extension = os.path.splitext(file_path)
            # if 1==1:
            if extension.lower() in file_types:
                convert_to_markdown(file_path, output_folder)
    
    msg_blurb = (f'See {output_folder}')
    msg_blurb_square(msg_blurb, color_green)    


def read_exif_data(file_path):
    """Reads selected EXIF data from an image file."""
    with open(file_path, 'rb') as f:
        tags = exifread.process_file(f, stop_tag='UNDEFINED')

    gps_latitude = tags.get('GPS GPSLatitude')
    gps_latitude_ref = tags.get('GPS GPSLatitudeRef')
    gps_longitude = tags.get('GPS GPSLongitude')
    gps_longitude_ref = tags.get('GPS GPSLongitudeRef')
    
    latitude = convert_to_decimal(gps_latitude.values, gps_latitude_ref.values[0]) if gps_latitude and gps_latitude_ref else ""
    longitude = convert_to_decimal(gps_longitude.values, gps_longitude_ref.values[0]) if gps_longitude and gps_longitude_ref else ""

    exif_data = {
        'Name': os.path.basename(file_path),
        'DateCreated': tags.get('Image DateTime'),
        'DateTimeOriginal': tags.get('EXIF DateTimeOriginal'),
        'FileCreateDate': datetime.fromtimestamp(os.path.getctime(file_path)).strftime('%Y-%m-%d %H:%M:%S'),
        'FileModifyDate': datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%Y-%m-%d %H:%M:%S'),
        'Timezone': tags.get('EXIF OffsetTime'),
        'DeviceManufacturer': tags.get('Image Make'),
        'LensMake': tags.get('EXIF LensMake'),
        'LensModel': tags.get('EXIF LensModel'),
        'Model': tags.get('Image Model'),
        'Software': tags.get('Image Software'),
        'LensInfo': tags.get('EXIF LensInfo'),
        'HostComputer': tags.get('Image HostComputer'),
        'FileSize': os.path.getsize(file_path),
        'FileType': os.path.splitext(file_path)[1].replace('.', '').upper(),
        'FileTypeExtension': os.path.splitext(file_path)[1].replace('.', '').lower(),
        'Altitude': tags.get('GPS GPSAltitude'),
        'Latitude': latitude,
        'Longitude': longitude,
        'Coordinate': f"{latitude},{longitude}" if latitude and longitude else "",
        'NumberOfImages': tags.get('Exif NumberOfImages'),
        'ExifToolVersion': tags.get('Exif ExifToolVersion'),
        'Icon': 'Images',
        'Type': 'Images',
        'Time': tags.get('EXIF DateTimeOriginal')
    }

    for key, value in exif_data.items():
        exif_data[key] = str(value) if value else ""
        
    Description = cleanup_description(exif_data)
    return exif_data, Description
    

def read_markdown_files(input_folder,output_folder):
    # Ensure input folder exists

    if not os.path.exists(output_folder):
        msg_blurb_square(f"Error: Output folder '{output_folder}' does not exist.", color_red)
        # msg_blurb = (f"Error: Output folder '{output_folder}' does not exist.")
        # msg_blurb_square(msg_blurb, color_red) 
        exit()

    if not os.path.exists(input_folder):
        msg_blurb_square(f"Error: Input folder '{input_folder}' does not exist.", color_red)
        # msg_blurb = (f"Error: Input folder '{input_folder}' does not exist.")
        # msg_blurb_square(msg_blurb, color_red) 
        exit()

    else:
        msg_blurb = (f"Reading markdown files from {input_folder}")
        msg_blurb_square(msg_blurb, color_green) 

    # List all markdown files (case-insensitive) containing "todo" in the filename
    md_files = [
        f for f in os.listdir(input_folder)
        if f.lower().endswith(".md")
    ]

    if not md_files:
        msg_blurb = (f"No markdown files found in {input_folder}")
        msg_blurb_square(msg_blurb, color_red) 
    else:
        for filename in md_files:
            md2pdf(filename, input_folder, output_folder)

    return md_files

    
def setup_obsidian(input_folder, output_folder):
    # Ensure the output folder exists, or create it if it doesn’t
    if not os.path.exists(output_folder):
        msg_blurb = f"The output folder '{output_folder}' doesn't exist. Would you like to create it? (y/n): "
        user_response = input(msg_blurb).strip().lower()

        if user_response == 'y':
            try:
                os.makedirs(output_folder, exist_ok=True)
                print(f"Output folder '{output_folder}' has been created.")
            except Exception as e:
                print(f"Error: Could not create the output folder. {e}")
                sys.exit(1)  # Exit with error code 1
        else:
            print("Exiting as output folder creation was declined.")
            sys.exit(1)  # Exit with error code 1

    msg_blurb = (f'See {output_folder}')
    msg_blurb_square(msg_blurb, color_green)

    app = {
      "alwaysUpdateLinks": True,
      "newFileLocation": "current",
      "newLinkFormat": "relative",
      "showUnsupportedFiles": True,
      "attachmentFolderPath": "Assets"
    }

    appearance = {
        "theme": "obsidian",
        "accentColor": "#5cb2f5",
        "interfaceFontFamily": "Times New Roman",
        "textFontFamily": "Times New Roman",
        "monospaceFontFamily": "Times New Roman",
        "baseFontSize": 17,
        "baseFontSizeAction": True
    }

    community_plugins = [
      "obsidian-excalidraw-plugin",
      "dataview",
      "obsidian-icon-folder",
      "obsidian-kanban",
      "calendar",
      "obsidian-tasks-plugin",
      "obsidian-advanced-slides",
      "obsidian-annotator",
      "homepage",
      "omnisearch",
      "periodic-notes",
      "url-into-selection",
      "obsidian-textgenerator-plugin",
      "obsidian-banners",
      "obsidian-plugin-toc",
      "media-extended",
      "settings-search",
      "emoji-shortcodes",
      "smart-connections",
      "obsidian-plugin-update-tracker",
      "janitor"
    ]

    core_plugins = {
        "file-explorer": True,
        "global-search": True,
        "switcher": True,
        "graph": True,
        "backlink": True,
        "canvas": True,
        "outgoing-link": True,
        "tag-pane": True,
        "properties": False,
        "page-preview": True,
        "daily-notes": True,
        "templates": True,
        "note-composer": True,
        "command-palette": True,
        "slash-command": False,
        "editor-status": True,
        "bookmarks": True,
        "markdown-importer": False,
        "zk-prefixer": False,
        "random-note": False,
        "outline": True,
        "word-count": True,
        "slides": False,
        "audio-recorder": False,
        "workspaces": False,
        "file-recovery": True,
        "publish": False,
        "sync": False
    }

    templates = {
      "folder": "assets/templates",
      "dateFormat": "YYYY-MM-DD",
      "timeFormat": "HH:mm:ss"
    }


    # Path to the .obsidian folder and the appearance.json file
    obsidian_folder = os.path.join(output_folder, '.obsidian')
    global assets_folder
    assets_folder = os.path.join(output_folder, '9.Assets')
    global assets_folder2
    assets_folder2 = os.path.join('9.Assets')
    global docs_folder
    docs_folder = os.path.join(output_folder, '4.Documents\originalDocuments')
    global media_folder
    # media_folder = os.path.join(output_folder, '9.Assets\media')
    media_folder = os.path.join(output_folder, (f'{assets_folder2}\media'))    
    global images_folder

    # images_folder = os.path.join(output_folder, (f'{assets_folder2}\media')) 

    images_folder = os.path.join(output_folder, '9.Assets\\media\\images')
    global scripts_folder
    scripts_folder = os.path.join(output_folder, '9.Assets\scripts')
    global templates_folder
    templates_folder = os.path.join(output_folder, '5.Templates')
    # templates_folder = os.path.join(output_folder, 'Template_2ndBrain')

    app_file = os.path.join(obsidian_folder, 'app.json')
    appearance_file = os.path.join(obsidian_folder, 'appearance.json')
    community_plugins_file = os.path.join(obsidian_folder, 'community-plugins.json')
    core_plugins_file = os.path.join(obsidian_folder, 'core-plugins.json')
    template_file = os.path.join(obsidian_folder, 'templates.json')

    # Check if the .obsidian folder exists, if not, create it
    if not os.path.exists(obsidian_folder):
        os.makedirs(obsidian_folder)
        print(f"Created {obsidian_folder}")

    # Check if the assets folder exists, if not, create it
    if not os.path.exists(assets_folder):
        os.makedirs(assets_folder)
        print(f"Created {assets_folder}")

    if not os.path.exists(media_folder):
        os.makedirs(media_folder)
        print(f"Created {media_folder}")

    # Check if the templates folder exists, if not, create it
    if not os.path.exists(templates_folder):
        os.makedirs(templates_folder)
        print(f"Created {templates_folder}")

    # Check if the Images folder exists, if not, create it
    if not os.path.exists(images_folder):
        os.makedirs(images_folder)
        print(f"Created {images_folder}")

    # Check if the documents folder exists, if not, create it
    if not os.path.exists(docs_folder):
        os.makedirs(docs_folder)
        print(f"Created {docs_folder}")

    # if not os.path.exists(templates_folder):
        # os.makedirs(templates_folder)
        # print(f"Created {templates_folder}")


    # Check if the scripts folder exists, if not, create it
    if not os.path.exists(scripts_folder):
        os.makedirs(scripts_folder)
        print(f"Created {scripts_folder}")


    if not os.path.exists(app_file):
        with open(app_file, 'w') as f:
            json.dump(app, f, indent=4)
        print(f"Created {app_file}")
    if not os.path.exists(appearance_file):
        with open(appearance_file, 'w') as f:
            json.dump(appearance, f, indent=4)
        print(f"Created {appearance_file}")
    if not os.path.exists(community_plugins_file):
        with open(community_plugins_file, 'w') as f:
            json.dump(community_plugins, f, indent=4)
        print(f"Created {community_plugins_file}")

    if not os.path.exists(core_plugins_file):
        with open(core_plugins_file, 'w') as f:
            json.dump(core_plugins, f, indent=4)
        print(f"Created {core_plugins_file}")


def Usage():
    file = sys.argv[0].split('\\')[-1]

    print(f'\nDescription: {color_green}{description}{color_reset}')
    print(f'{file} Version: {version} by {author}')
    print(f"    {file} -b")
    print(f"    {file} -c")
    print(f"    {file} -c -I C:\Forensics\scripts\python\Files -O ObsidianNotebook") 
    print(f"    {file} -c -O ObsidianNotebook") 
    print(f"    {file} -c -I test_files -O ObsidianNotebook") 
    print(f"    python {file} -P -I C:\Forensics\ObsidianVaults\ForensicsBrain\8.Tasks -O C:\Forensics\scripts\python\playfolder")     
    print(f'    python {file} -P -I \"B:\Agent Folders\Forensic Examiners\B-drive_SecondBrain\8.Tasks" -O \"B:\Administrative\DFE_Folder\"')     


    print(f"    {file} -b -t -T Template_Cases") 


if __name__ == '__main__':
    main()


# <<<<<<<<<<<<<<<<<<<<<<<<<< Revision History >>>>>>>>>>>>>>>>>>>>>>>>>>


"""
1.1.2 -P option spits a backup .txt file if the todo is longer than 50 lines
1.1.1 - markdown files to pdfs -C
0.2.1 - .DOCX and .PPTX conversions with metadata
0.1.5 - add exifdata to the bottom of the .md file
0.1.4 - (creation_time, access_time, modified_time) = get_file_timestamps(file_path)
0.1.3 - -t option to create a default obsidian folder/files
0.1.2 - create obsidian config files if they don't exist
0.1.1 - Convert the content of .txt, .pdf, .docx, and .eml to Markdown, for use in Obsidian
0.0.9 - converted to template version
0.0.1 - created by ChatGPT
"""


# <<<<<<<<<<<<<<<<<<<<<<<<<< Future Wishlist  >>>>>>>>>>>>>>>>>>>>>>>>>>

"""

can i do -C to a template pdf? Can I add a graphic? Can I add bold formatting?
export to html -w

pdf convert data and tables.
undo the try except in .eml multi threaded

fix rtf converter - doesn't do all formatting like bold

The script is Windows-specific checks (e.g., file creation time) could potentially be improved for cross-platform compatibility. You might use pathlib for better handling of file paths across platforms.

this will overwrite files with the same name, create a way of if file exists, save it with a unique name

"""

# <<<<<<<<<<<<<<<<<<<<<<<<<<      notes            >>>>>>>>>>>>>>>>>>>>>>>>>>

"""
# obsidianShortcut.cmd
start "" "C:\Program Files\Obsidian\Obsidian.exe" "C:\Forensics\ObsidianVaults\ForensicsBrain"

"""

# <<<<<<<<<<<<<<<<<<<<<<<<<<      The End        >>>>>>>>>>>>>>>>>>>>>>>>>>
